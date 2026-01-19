# -*- coding: utf-8 -*-
"""
Ferramentas de Geração de Documentos para o Agente IA
Excel, Word, PowerPoint, PDF, Markdown
"""

import io
import logging
from typing import Optional, Dict, Any, List
from datetime import datetime

import pandas as pd

logger = logging.getLogger(__name__)

# ============================================================================
# EXCEL
# ============================================================================

def gerar_excel(
    dados: pd.DataFrame,
    nome_planilha: str = "Dados",
    titulo: str = None,
    formatar: bool = True
) -> bytes:
    """
    Gera arquivo Excel formatado.
    
    Args:
        dados: DataFrame com os dados
        nome_planilha: Nome da aba
        titulo: Título opcional na primeira linha
        formatar: Se deve aplicar formatação
    
    Returns:
        Bytes do arquivo Excel
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils.dataframe import dataframe_to_rows
    from openpyxl.utils import get_column_letter
    
    wb = Workbook()
    ws = wb.active
    ws.title = nome_planilha
    
    start_row = 1
    
    # Título se fornecido
    if titulo:
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(dados.columns))
        cell = ws.cell(row=1, column=1, value=titulo)
        cell.font = Font(bold=True, size=14, color="7C3AED")
        cell.alignment = Alignment(horizontal='center')
        start_row = 3
    
    # Escrever dados
    for r_idx, row in enumerate(dataframe_to_rows(dados, index=False, header=True), start_row):
        for c_idx, value in enumerate(row, 1):
            cell = ws.cell(row=r_idx, column=c_idx, value=value)
            
            if formatar:
                # Header
                if r_idx == start_row:
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill(start_color="7C3AED", end_color="7C3AED", fill_type="solid")
                    cell.alignment = Alignment(horizontal='center')
                
                # Bordas
                thin_border = Border(
                    left=Side(style='thin', color='CCCCCC'),
                    right=Side(style='thin', color='CCCCCC'),
                    top=Side(style='thin', color='CCCCCC'),
                    bottom=Side(style='thin', color='CCCCCC')
                )
                cell.border = thin_border
                
                # Números
                if isinstance(value, (int, float)) and r_idx > start_row:
                    if value > 1000:
                        cell.number_format = '#,##0.00'
                    elif value < 1:
                        cell.number_format = '0.00%'
    
    # Ajustar largura das colunas
    for col_idx, column_cells in enumerate(ws.columns, 1):
        max_length = 0
        column = get_column_letter(col_idx)
        for cell in column_cells:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        ws.column_dimensions[column].width = min(max_length + 2, 50)
    
    # Congelar painel no header
    ws.freeze_panes = ws.cell(row=start_row + 1, column=1)
    
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    
    logger.info(f"Excel gerado: {len(dados)} linhas, {len(dados.columns)} colunas")
    return buf.getvalue()


def ler_excel(arquivo: bytes, planilha: str = None) -> pd.DataFrame:
    """
    Lê arquivo Excel e retorna DataFrame.
    
    Args:
        arquivo: Bytes do arquivo Excel
        planilha: Nome da planilha (opcional, usa a primeira)
    
    Returns:
        DataFrame com os dados
    """
    buf = io.BytesIO(arquivo)
    df = pd.read_excel(buf, sheet_name=planilha)
    logger.info(f"Excel lido: {len(df)} linhas, {len(df.columns)} colunas")
    return df


# ============================================================================
# WORD (DOCX)
# ============================================================================

def gerar_word(
    titulo: str,
    secoes: List[Dict[str, Any]],
    autor: str = "Agente IA - Propensão"
) -> bytes:
    """
    Gera documento Word formatado.
    
    Args:
        titulo: Título do documento
        secoes: Lista de seções com 'titulo', 'conteudo', 'tipo' (texto, tabela, imagem)
        autor: Nome do autor
    
    Returns:
        Bytes do arquivo .docx
    """
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    
    doc = Document()
    
    # Configurar estilos
    style = doc.styles['Title']
    style.font.color.rgb = RGBColor(124, 58, 237)  # Purple
    
    # Título principal
    title = doc.add_heading(titulo, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Metadados
    doc.add_paragraph(f"Gerado em: {datetime.now().strftime('%d/%m/%Y às %H:%M')}")
    doc.add_paragraph(f"Autor: {autor}")
    doc.add_paragraph("")
    
    # Processar seções
    for secao in secoes:
        tipo = secao.get("tipo", "texto")
        
        if secao.get("titulo"):
            doc.add_heading(secao["titulo"], level=1)
        
        if tipo == "texto":
            for paragrafo in secao.get("conteudo", "").split("\n\n"):
                if paragrafo.strip():
                    doc.add_paragraph(paragrafo.strip())
        
        elif tipo == "tabela":
            dados = secao.get("dados")
            if isinstance(dados, pd.DataFrame):
                # Criar tabela
                table = doc.add_table(rows=1, cols=len(dados.columns))
                table.style = 'Table Grid'
                
                # Header
                hdr_cells = table.rows[0].cells
                for i, col in enumerate(dados.columns):
                    hdr_cells[i].text = str(col)
                    hdr_cells[i].paragraphs[0].runs[0].font.bold = True
                
                # Dados
                for _, row in dados.iterrows():
                    row_cells = table.add_row().cells
                    for i, value in enumerate(row):
                        row_cells[i].text = str(value)
        
        elif tipo == "imagem":
            imagem_bytes = secao.get("imagem")
            if imagem_bytes:
                img_stream = io.BytesIO(imagem_bytes)
                doc.add_picture(img_stream, width=Inches(6))
                last_para = doc.paragraphs[-1]
                last_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        elif tipo == "lista":
            for item in secao.get("itens", []):
                doc.add_paragraph(item, style='List Bullet')
    
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    
    logger.info(f"Word gerado: {len(secoes)} seções")
    return buf.getvalue()


# ============================================================================
# POWERPOINT (PPTX)
# ============================================================================

def gerar_powerpoint(
    titulo: str,
    slides: List[Dict[str, Any]],
    autor: str = "Agente IA - Propensão"
) -> bytes:
    """
    Gera apresentação PowerPoint.
    
    Args:
        titulo: Título da apresentação
        slides: Lista de slides com 'titulo', 'conteudo', 'tipo' (texto, bullets, imagem, tabela)
        autor: Nome do autor
    
    Returns:
        Bytes do arquivo .pptx
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide de título
    title_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Background escuro
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(10, 10, 10)
    background.line.fill.background()
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = titulo
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(168, 85, 247)  # Purple
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"{autor} | {datetime.now().strftime('%d/%m/%Y')}"
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = RGBColor(148, 163, 184)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Processar slides de conteúdo
    for slide_data in slides:
        content_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(content_layout)
        
        # Background
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(10, 10, 10)
        bg.line.fill.background()
        
        # Título do slide
        if slide_data.get("titulo"):
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data["titulo"]
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        tipo = slide_data.get("tipo", "texto")
        
        if tipo == "texto":
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for para_text in slide_data.get("conteudo", "").split("\n"):
                para = content_frame.add_paragraph()
                para.text = para_text
                para.font.size = Pt(18)
                para.font.color.rgb = RGBColor(226, 232, 240)
        
        elif tipo == "bullets":
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
            content_frame = content_box.text_frame
            
            for item in slide_data.get("itens", []):
                para = content_frame.add_paragraph()
                para.text = f"• {item}"
                para.font.size = Pt(20)
                para.font.color.rgb = RGBColor(226, 232, 240)
                para.space_after = Pt(12)
        
        elif tipo == "imagem":
            imagem_bytes = slide_data.get("imagem")
            if imagem_bytes:
                img_stream = io.BytesIO(imagem_bytes)
                slide.shapes.add_picture(img_stream, Inches(1), Inches(1.3), width=Inches(11.333))
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    
    logger.info(f"PowerPoint gerado: {len(slides) + 1} slides")
    return buf.getvalue()


# ============================================================================
# PDF (com ReportLab)
# ============================================================================

def gerar_pdf(
    titulo: str,
    secoes: List[Dict[str, Any]],
    autor: str = "Agente IA - Propensão",
    tema: str = "dark"
) -> bytes:
    """
    Gera relatório PDF profissional.
    
    Args:
        titulo: Título do relatório
        secoes: Lista de seções com conteúdo
        autor: Nome do autor
        tema: 'dark' ou 'light'
    
    Returns:
        Bytes do arquivo PDF
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    
    # Cores do tema
    if tema == "dark":
        bg_color = colors.HexColor("#0a0a0a")
        text_color = colors.white
        accent_color = colors.HexColor("#a855f7")
    else:
        bg_color = colors.white
        text_color = colors.HexColor("#1e293b")
        accent_color = colors.HexColor("#7c3aed")
    
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=2*cm, bottomMargin=2*cm)
    
    # Estilos
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name='TituloCustom',
        parent=styles['Title'],
        fontSize=24,
        textColor=accent_color,
        spaceAfter=30,
        alignment=TA_CENTER
    ))
    styles.add(ParagraphStyle(
        name='SubtituloCustom',
        parent=styles['Heading1'],
        fontSize=16,
        textColor=accent_color,
        spaceBefore=20,
        spaceAfter=10
    ))
    styles.add(ParagraphStyle(
        name='CorpoCustom',
        parent=styles['Normal'],
        fontSize=11,
        textColor=text_color if tema == "light" else colors.HexColor("#94a3b8"),
        spaceAfter=12
    ))
    
    elementos = []
    
    # Título
    elementos.append(Paragraph(titulo, styles['TituloCustom']))
    elementos.append(Paragraph(
        f"Gerado em {datetime.now().strftime('%d/%m/%Y às %H:%M')} | {autor}",
        styles['CorpoCustom']
    ))
    elementos.append(Spacer(1, 30))
    
    # Processar seções
    for secao in secoes:
        if secao.get("titulo"):
            elementos.append(Paragraph(secao["titulo"], styles['SubtituloCustom']))
        
        tipo = secao.get("tipo", "texto")
        
        if tipo == "texto":
            for para in secao.get("conteudo", "").split("\n\n"):
                if para.strip():
                    elementos.append(Paragraph(para.strip(), styles['CorpoCustom']))
        
        elif tipo == "tabela":
            dados = secao.get("dados")
            if isinstance(dados, pd.DataFrame):
                table_data = [dados.columns.tolist()] + dados.values.tolist()
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), accent_color),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('GRID', (0, 0), (-1, -1), 1, colors.HexColor("#cccccc")),
                    ('FONTSIZE', (0, 1), (-1, -1), 9),
                ]))
                elementos.append(t)
                elementos.append(Spacer(1, 15))
        
        elif tipo == "imagem":
            imagem_bytes = secao.get("imagem")
            if imagem_bytes:
                img_stream = io.BytesIO(imagem_bytes)
                img = Image(img_stream, width=16*cm, height=10*cm)
                elementos.append(img)
                elementos.append(Spacer(1, 15))
        
        elif tipo == "lista":
            for item in secao.get("itens", []):
                elementos.append(Paragraph(f"• {item}", styles['CorpoCustom']))
    
    doc.build(elementos)
    buf.seek(0)
    
    logger.info(f"PDF gerado: {len(secoes)} seções")
    return buf.getvalue()


# ============================================================================
# MARKDOWN
# ============================================================================

def gerar_markdown(
    titulo: str,
    secoes: List[Dict[str, Any]],
    autor: str = "Agente IA - Propensão"
) -> bytes:
    """
    Gera documento Markdown.
    
    Args:
        titulo: Título do documento
        secoes: Lista de seções
        autor: Nome do autor
    
    Returns:
        Bytes do arquivo .md
    """
    linhas = [
        f"# {titulo}",
        "",
        f"*Gerado em {datetime.now().strftime('%d/%m/%Y às %H:%M')} por {autor}*",
        "",
        "---",
        ""
    ]
    
    for secao in secoes:
        if secao.get("titulo"):
            nivel = secao.get("nivel", 2)
            linhas.append(f"{'#' * nivel} {secao['titulo']}")
            linhas.append("")
        
        tipo = secao.get("tipo", "texto")
        
        if tipo == "texto":
            linhas.append(secao.get("conteudo", ""))
            linhas.append("")
        
        elif tipo == "tabela":
            dados = secao.get("dados")
            if isinstance(dados, pd.DataFrame):
                linhas.append(dados.to_markdown(index=False))
                linhas.append("")
        
        elif tipo == "lista":
            for item in secao.get("itens", []):
                linhas.append(f"- {item}")
            linhas.append("")
        
        elif tipo == "codigo":
            linguagem = secao.get("linguagem", "")
            codigo = secao.get("codigo", "")
            linhas.append(f"```{linguagem}")
            linhas.append(codigo)
            linhas.append("```")
            linhas.append("")
    
    conteudo = "\n".join(linhas)
    
    logger.info(f"Markdown gerado: {len(secoes)} seções")
    return conteudo.encode('utf-8')


__all__ = [
    "gerar_excel",
    "ler_excel",
    "gerar_word",
    "gerar_powerpoint",
    "gerar_pdf",
    "gerar_markdown"
]
